# -*- coding: utf-8 -*-
"""RAG.ipynb

Automatically generated by Colab.

Original file is located at
    https://colab.research.google.com/drive/11TK1xmybI7ja8j33JoSkcs_m4jHax7fA
"""

# !pip install langchain langchain_openai chromadb streamlit unstructured python-pptx python-docx pypdf google-auth-oauthlib google-auth-httplib2 google-api-python-client langchain_community PyMuPDF docx2txt google-cloud-storage google.cloud tenacity



import os
import tempfile
import streamlit as st
import io
import base64
from typing import List, Tuple
import fitz
import docx2txt
from tenacity import retry, stop_after_attempt, wait_exponential
from google.cloud import storage

from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain.retrievers import ContextualCompressionRetriever
from langchain.retrievers.document_compressors import LLMChainExtractor
from langchain.schema import Document
from langchain.chains import ConversationalRetrievalChain

import pypdf
import json
from pptx import Presentation
from docx import Document

BUCKET_NAME = os.environ.get("BUCKET_NAME")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
PRIVATE_KEY_ID = os.environ.get("PRIVATE_KEY_ID")
PRIVATE_KEY = os.environ.get("PRIVATE_KEY")

import streamlit
import json
    service_account_info = {
        "PRIVATE_KEY_ID": private_key_id,
        "PRIVATE_KEY": private_key,
    }

def initialize_gcs_client():
    """Initialize Google Cloud Storage client with service account"""
    storage_client = storage.Client.from_service_account_json('service_account_info')
    return storage_client

def download_files_from_bucket(storage_client) -> List[Tuple[str, bytes]]:
    """Download all files from the specified GCS bucket"""
    bucket = storage_client.bucket(BUCKET_NAME)
    blobs = bucket.list_blobs()

    downloaded_files = []
    for blob in blobs:
        content = blob.download_as_bytes()
        downloaded_files.append((blob.name, content))
        print(f"Downloaded: {blob.name}")

    return downloaded_files

def process_file(file_name: str, content: bytes) -> str:
    """Process different file types and extract text"""
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        temp_file.write(content)
        temp_path = temp_file.name

    text = ""
    try:
        if file_name.endswith('.pdf'):
            # Use PyMuPDF for PDFs
            doc = fitz.open(temp_path)
            for page in doc:
                text += page.get_text()
            doc.close()

        elif file_name.endswith('.docx'):
            # Use docx2txt for DOCX files
            text = docx2txt.process(temp_path)

        elif file_name.endswith('.pptx'):
            prs = Presentation(temp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

    finally:
        os.unlink(temp_path)

    return text

def initialize_qa_system():
    """Initialize the QA system with the configured bucket and API key"""
    # Initialize OpenAI components
    embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
    llm = ChatOpenAI(
        temperature=0,
        openai_api_key=OPENAI_API_KEY,
        model="gpt-4-turbo-1106"
    )

    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=1, max=10))
    def embed_with_retry(texts):
        return embeddings.embed_documents(texts)

    # Download and process files from GCS
    storage_client = initialize_gcs_client()
    files = download_files_from_bucket(storage_client)

    print(f"Files found: {[f[0] for f in files]}")

    all_text = ""
    for file_name, content in files:
        print(f"Processing: {file_name}")
        all_text += process_file(file_name, content) + "\n\n"

    print(f"All text length: {len(all_text)}")

    # Split text into chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.create_documents([all_text])

    # Create vector store
    vectorstore = Chroma.from_documents(
        documents=texts,
        embedding=embeddings,
        persist_directory="chroma_db"
    )

    # Set up retriever with reranking
    base_retriever = vectorstore.as_retriever(search_kwargs={"k": 10})
    compressor = LLMChainExtractor.from_llm(llm)
    compression_retriever = ContextualCompressionRetriever(
        base_compressor=compressor,
        base_retriever=base_retriever
    )

    # Create conversational chain
    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=compression_retriever,
        return_source_documents=True,
        verbose=True
    )

    return qa_chain

def add_logo():
    """Add logo to the Streamlit interface"""
    logo_svg = '''
    <svg width="100" height="100" viewBox="0 0 100 100" xmlns="http://www.w3.org/2000/svg">
        <rect x="20" y="20" width="60" height="60" fill="#4A90E2"/>
        <circle cx="50" cy="50" r="25" fill="white"/>
        <text x="50" y="55" font-family="Arial" font-size="14" fill="#4A90E2" text-anchor="middle">RAG</text>
    </svg>
    '''

    b64 = base64.b64encode(logo_svg.encode('utf-8')).decode()

    st.markdown(
        f"""
        <style>
            [data-testid="stHeader"] {{
                background-color: white;
            }}

            .logo-container {{
                display: flex;
                justify-content: center;
                margin-bottom: 2rem;
            }}

            .chat-container {{
                max-width: 800px;
                margin: 0 auto;
                padding: 2rem;
            }}

            .stChatMessage {{
                background-color: #ffffff;
                border-radius: 10px;
                padding: 1rem;
                margin: 0.5rem 0;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }}
        </style>
        <div class="logo-container">
            <img src="data:image/svg+xml;base64,{b64}" alt="Logo" width="100">
        </div>
        """,
        unsafe_allow_html=True
    )

def main():
    """Main Streamlit application"""
    st.set_page_config(page_title="Document Q&A System", layout="wide")
    add_logo()

    st.markdown('<h1 style="text-align: center;">Ask me anything about the fund</h1>', unsafe_allow_html=True)

    # Initialize the QA system if not already done
    if "qa_chain" not in st.session_state:
        with st.spinner("Initializing the system..."):
            st.session_state.qa_chain = initialize_qa_system()
            st.session_state.chat_history = []

    # Chat interface
    st.markdown('<div class="chat-container">', unsafe_allow_html=True)

    if "messages" not in st.session_state:
        st.session_state.messages = []

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if prompt := st.chat_input("What would you like to know?"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        with st.chat_message("assistant"):
            response = st.session_state.qa_chain(
                {"question": prompt, "chat_history": st.session_state.chat_history}
            )
            st.markdown(response["answer"])
            st.session_state.messages.append({"role": "assistant", "content": response["answer"]})
            st.session_state.chat_history.append((prompt, response["answer"]))

    st.markdown('</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
